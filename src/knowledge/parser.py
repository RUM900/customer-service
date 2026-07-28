"""
多格式文档解析器

支持: PDF, DOCX, MD, TXT, CSV, HTML, PPTX
每种格式返回统一的 ParsedDocument。
"""
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    """解析后的文档"""
    title: str = ""
    content: str = ""                # 纯文本内容
    source: str = ""                 # 来源文件名
    format: str = ""                 # 原始格式
    metadata: dict = field(default_factory=dict)
    page_count: int = 0


def parse_file(file_path: str, original_filename: str = "") -> ParsedDocument:
    """
    自动检测格式并解析文件

    Args:
        file_path: 文件在磁盘上的路径（临时文件）
        original_filename: 原始文件名（用于检测格式）

    Returns:
        ParsedDocument
    """
    path = Path(file_path)
    filename = original_filename or path.name
    suffix = Path(filename).suffix.lower()

    parsers = {
        ".pdf":  _parse_pdf,
        ".docx": _parse_docx,
        ".md":   _parse_markdown,
        ".txt":  _parse_text,
        ".csv":  _parse_csv,
        ".html": _parse_html,
        ".htm":  _parse_html,
        ".pptx": _parse_pptx,
    }

    parser = parsers.get(suffix)
    if parser is None:
        raise ValueError(f"不支持的文件格式: {suffix}。支持: {', '.join(parsers.keys())}")

    logger.info(f"解析文档: {filename} ({suffix})")
    return parser(str(path), filename)


# ============================================================
# PDF
# ============================================================

def _parse_pdf(file_path: str, filename: str) -> ParsedDocument:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)

        return ParsedDocument(
            title=filename,
            content="\n\n".join(pages),
            source=filename,
            format="pdf",
            page_count=len(reader.pages),
            metadata={"filename": filename},
        )
    except ImportError:
        raise ImportError("解析 PDF 需要 PyPDF2: pip install PyPDF2")


# ============================================================
# DOCX
# ============================================================

def _parse_docx(file_path: str, filename: str) -> ParsedDocument:
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                if row_text.strip():
                    paragraphs.append(row_text)

        return ParsedDocument(
            title=filename,
            content="\n\n".join(paragraphs),
            source=filename,
            format="docx",
            metadata={"filename": filename},
        )
    except ImportError:
        raise ImportError("解析 DOCX 需要 python-docx: pip install python-docx")


# ============================================================
# Markdown
# ============================================================

def _parse_markdown(file_path: str, filename: str) -> ParsedDocument:
    with open(file_path, "r", encoding="utf-8") as f:
        raw = f.read()

    # 提取标题（第一个 # 开头行）
    title = filename
    for line in raw.split("\n"):
        if line.startswith("# "):
            title = line[2:].strip()
            break

    return ParsedDocument(
        title=title,
        content=raw,
        source=filename,
        format="markdown",
        metadata={"filename": filename},
    )


# ============================================================
# TXT
# ============================================================

def _parse_text(file_path: str, filename: str) -> ParsedDocument:
    # 尝试 UTF-8，失败则 GBK
    for enc in ["utf-8", "gbk", "latin-1"]:
        try:
            with open(file_path, "r", encoding=enc) as f:
                content = f.read()
            break
        except UnicodeDecodeError:
            continue
    else:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

    return ParsedDocument(
        title=filename,
        content=content,
        source=filename,
        format="txt",
        metadata={"filename": filename},
    )


# ============================================================
# CSV
# ============================================================

def _parse_csv(file_path: str, filename: str) -> ParsedDocument:
    import csv
    import io

    with open(file_path, "r", encoding="utf-8-sig", errors="ignore") as f:
        content = f.read()

    # 将 CSV 转为可读文本：每行 "列名: 值, 列名: 值"
    reader = csv.DictReader(io.StringIO(content))
    rows = []
    for i, row in enumerate(reader):
        parts = [f"{k}: {v}" for k, v in row.items() if v]
        if parts:
            rows.append(f"[行{i+1}] " + " | ".join(parts))

    return ParsedDocument(
        title=filename,
        content="\n".join(rows),
        source=filename,
        format="csv",
        metadata={"filename": filename},
    )


# ============================================================
# HTML
# ============================================================

def _parse_html(file_path: str, filename: str) -> ParsedDocument:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()

    # 简单去标签
    try:
        from html.parser import HTMLParser

        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.texts = []

            def handle_data(self, data):
                text = data.strip()
                if text:
                    self.texts.append(text)

        extractor = TextExtractor()
        extractor.feed(raw)
        content = "\n".join(extractor.texts)
    except Exception:
        # 正则兜底
        import re
        content = re.sub(r"<[^>]+>", " ", raw)
        content = re.sub(r"\s+", " ", content).strip()

    # 提取 title
    title = filename
    import re
    match = re.search(r"<title>(.*?)</title>", raw, re.IGNORECASE)
    if match:
        title = match.group(1)

    return ParsedDocument(
        title=title,
        content=content,
        source=filename,
        format="html",
        metadata={"filename": filename},
    )


# ============================================================
# PPTX
# ============================================================

def _parse_pptx(file_path: str, filename: str) -> ParsedDocument:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))

        return ParsedDocument(
            title=filename,
            content="\n\n".join(slides),
            source=filename,
            format="pptx",
            page_count=len(prs.slides),
            metadata={"filename": filename},
        )
    except ImportError:
        raise ImportError("解析 PPTX 需要 python-pptx: pip install python-pptx")
